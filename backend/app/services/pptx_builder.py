"""
pptx_builder.py — Schumacher Homes Monthly Report PPTX

Matches the Single Grain presentation style:
  • Warm off-white background (#F2F0EB)
  • Orange/rust accent (#C4622D)
  • Dark charcoal text (#2D2926)
  • 70/30 split: left content + right orange panel
  • Single Grain logo bottom-left
  • Georgia serif headings, Calibri sans-serif body
"""

from __future__ import annotations

import io
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree

# ─── Brand Palette ────────────────────────────────────────────────────────────
BG       = RGBColor(0xF2, 0xF0, 0xEB)   # warm off-white slide background
ORANGE   = RGBColor(0xC4, 0x62, 0x2D)   # primary orange/rust accent
DARK     = RGBColor(0x2D, 0x29, 0x26)   # near-black charcoal text
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREEN    = RGBColor(0x2D, 0x7A, 0x3A)   # positive MoM / top performer
RED      = RGBColor(0xB9, 0x1C, 0x1C)   # negative MoM / needs attention
TEAL_G   = RGBColor(0x15, 0x60, 0x3A)   # "On Track" / "Healthy" status
ROW_ALT  = RGBColor(0xEE, 0xEB, 0xE6)   # alternate table row background
PEACH    = RGBColor(0xED, 0xC8, 0xB0)   # decorative circle (slides 1 & 6)
GRAY_SUB = RGBColor(0x78, 0x74, 0x70)   # subtitle / metadata text
GRAY_BDR = RGBColor(0xD4, 0xD0, 0xCA)   # card border

# ─── Layout Constants ─────────────────────────────────────────────────────────
SW = Inches(13.333)   # slide width  (16:9 widescreen)
SH = Inches(7.5)      # slide height
RP = Inches(9.2)      # right-panel left edge  (~69 % of width)
RW = SW - RP          # right-panel width  (~4.13 ")

# ─── Utility helpers ──────────────────────────────────────────────────────────

def _to_list(val) -> List[str]:
    """Normalise a value (string or list) to a list of non-empty strings."""
    if isinstance(val, list):
        out = []
        for x in val:
            s = str(x).strip().lstrip("-").lstrip("•").lstrip("●").strip()
            if s:
                out.append(s)
        return out
    if isinstance(val, str):
        out = []
        for line in val.split("\n"):
            line = line.strip().lstrip("-").lstrip("•").lstrip("●").strip()
            if line:
                out.append(line)
        return out
    return []


def _get_slide_content(report: Dict, slide_number: int) -> Dict:
    """Return the content dict for the requested slide number."""
    for s in report.get("slides", []):
        if s.get("slide_number") == slide_number:
            return s.get("content", {})
    return {}


def _status_color(status: str) -> RGBColor:
    s = str(status).lower()
    if any(x in s for x in ("on track", "healthy", "complete")):
        return TEAL_G
    if "monitor" in s:
        return ORANGE
    return ORANGE   # "in progress" etc.


# ─── Low-level drawing primitives ─────────────────────────────────────────────

def _blank_slide(prs: Presentation):
    """Add a blank slide and fill the background with BG colour."""
    blank = next(
        (l for l in prs.slide_layouts if l.name.lower() == "blank"),
        prs.slide_layouts[-1],
    )
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def _box(slide, l, t, w, h,
         fill_color: Optional[RGBColor] = None,
         line_color: Optional[RGBColor] = None,
         line_pt: float = 0.5):
    """Add a rectangle shape with optional fill and border."""
    shape = slide.shapes.add_shape(1, l, t, w, h)   # 1 = Rectangle
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def _oval(slide, l, t, w, h, fill_color: RGBColor):
    """Add a filled oval/circle."""
    shape = slide.shapes.add_shape(9, l, t, w, h)   # 9 = Oval
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _txt(slide, text: str, l, t, w, h,
         size: float = 11,
         bold: bool = False,
         italic: bool = False,
         color: RGBColor = DARK,
         align: PP_ALIGN = PP_ALIGN.LEFT,
         font: str = "Calibri",
         wrap: bool = True):
    """Add a single-run text box."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txb


# ─── Reusable slide components ────────────────────────────────────────────────

def _sg_logo(slide):
    """Single Grain logo — orange flame oval + 'SINGLE GRAIN' text."""
    _oval(slide, Inches(0.28), SH - Inches(0.52), Inches(0.3), Inches(0.38), ORANGE)
    _txt(slide, "SINGLE GRAIN",
         Inches(0.65), SH - Inches(0.44),
         Inches(1.6), Inches(0.3),
         size=7.5, color=DARK, font="Calibri")


def _stat_bar(slide, stats: List[tuple], x=None, y=None, w=None, h=None):
    """Orange stat bar.  stats = [(label, value), ...]"""
    x = x if x is not None else Inches(0.4)
    y = y if y is not None else Inches(1.1)
    w = w if w is not None else (RP - Inches(0.8))
    h = h if h is not None else Inches(0.82)
    _box(slide, x, y, w, h, fill_color=ORANGE)
    col_w = w / len(stats)
    for i, (label, value) in enumerate(stats):
        cx = x + col_w * i
        _txt(slide, label,
             cx + Inches(0.05), y + Inches(0.06),
             col_w - Inches(0.05), Inches(0.22),
             size=8.5, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
        _txt(slide, str(value),
             cx, y + Inches(0.28),
             col_w, Inches(0.46),
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")


def _slide_header(slide, title: str, subtitle: str = ""):
    """Large bold title + optional grey subtitle line."""
    _txt(slide, title,
         Inches(0.4), Inches(0.15),
         RP - Inches(0.5), Inches(0.65),
         size=28, bold=True, color=DARK, font="Georgia")
    if subtitle:
        _txt(slide, subtitle,
             Inches(0.4), Inches(0.75),
             RP - Inches(0.5), Inches(0.26),
             size=9.5, color=GRAY_SUB, font="Calibri")


def _right_panel(slide, title1: str, bullets1: List[str],
                 title2: str = "", bullets2: List[str] = None):
    """Draw the standard orange right panel with one or two bullet sections."""
    bullets2 = bullets2 or []
    _box(slide, RP, Inches(0), RW, SH, fill_color=ORANGE)

    y = Inches(0.42)

    def _section(header, items):
        nonlocal y
        _txt(slide, header,
             RP + Inches(0.18), y,
             RW - Inches(0.28), Inches(0.55),
             size=19, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font="Georgia")
        y += Inches(0.6)
        for item in items:
            _oval(slide,
                  RP + Inches(0.15), y + Inches(0.07),
                  Inches(0.09), Inches(0.09), WHITE)
            _txt(slide, item,
                 RP + Inches(0.32), y,
                 RW - Inches(0.46), Inches(0.95),
                 size=9.5, color=WHITE, font="Calibri", wrap=True)
            y += Inches(0.95)
        y += Inches(0.08)

    _section(title1, bullets1)
    if title2:
        _section(title2, bullets2)


# ─── Slide builders ───────────────────────────────────────────────────────────

def _slide1(prs: Presentation, report: Dict):
    """Slide 1 — Agenda."""
    slide = _blank_slide(prs)
    c = _get_slide_content(report, 1)
    month = report.get("report_month", "Monthly Report")

    # Large decorative peach circle — top right, partially off-screen
    _oval(slide, Inches(9.6), Inches(-1.8), Inches(5.8), Inches(5.8), PEACH)

    # Title
    _txt(slide, month,
         Inches(0.5), Inches(2.0), Inches(5.6), Inches(0.9),
         size=36, bold=True, color=DARK, font="Georgia")

    # "Agenda ●" — two-run textbox
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(2.88), Inches(5.6), Inches(0.9))
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Agenda "
    r1.font.size = Pt(40)
    r1.font.bold = True
    r1.font.color.rgb = DARK
    r1.font.name = "Georgia"
    r2 = p.add_run()
    r2.text = "●"
    r2.font.size = Pt(28)
    r2.font.bold = False
    r2.font.color.rgb = ORANGE
    r2.font.name = "Calibri"

    # Dashed divider line
    _box(slide, Inches(0.5), Inches(3.88), Inches(5.5), Pt(1.5), fill_color=GRAY_SUB)

    # Agenda items — right column with orange arrows
    agenda = c.get("agenda", [
        "Performance Reporting & Data Integrity.",
        "Ad Development & Testing.",
        "Priority Updates / New Business.",
    ])
    y = Inches(2.25)
    for item in agenda:
        _txt(slide, "➤",
             Inches(6.5), y, Inches(0.35), Inches(0.55),
             size=14, color=ORANGE, font="Calibri")
        _txt(slide, item,
             Inches(6.92), y, Inches(5.9), Inches(0.65),
             size=14, bold=True, color=DARK, font="Georgia", wrap=True)
        y += Inches(1.05)

    _sg_logo(slide)
    return slide


def _slide2(prs: Presentation, report: Dict):
    """Slide 2 — Paid Media KPIs & MoM Analysis."""
    slide = _blank_slide(prs)
    c = _get_slide_content(report, 2)
    month = report.get("report_month", "Monthly")

    prev_label  = c.get("prev_month_label", "Prior Month")
    curr_label  = c.get("curr_month_label", month)
    stats_raw   = c.get("summary_stats", {})
    total_spend = stats_raw.get("total_spend", 0)
    total_leads = stats_raw.get("total_leads", 0)
    blended_cpl = stats_raw.get("blended_cpl", 0)

    _slide_header(slide, "Paid Media KPIs & MoM Analysis",
                  f"Ad Platform Data  |  {curr_label}  |  32 Locations")

    _stat_bar(slide, [
        ("Total Spend", f"${total_spend:,.0f}"),
        ("Leads",       f"{total_leads:,}"),
        ("Avg CPL",     f"${blended_cpl:,.2f}"),
    ], y=Inches(1.08))

    # MoM table
    mom_rows = c.get("mom_table", [])
    if mom_rows:
        tx      = Inches(0.4)
        ty      = Inches(2.08)
        cw_list = [Inches(2.35), Inches(1.55), Inches(1.55), Inches(1.45), Inches(0.38)]
        row_h   = Inches(0.305)
        total_w = sum(cw_list)
        headers = ["Metric", prev_label, curr_label, "MoM Change", ""]

        _box(slide, tx, ty, total_w, row_h, fill_color=ORANGE)
        cx = tx
        for hdr, cw in zip(headers, cw_list):
            _txt(slide, hdr, cx + Inches(0.04), ty + Inches(0.04),
                 cw - Inches(0.04), row_h,
                 size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
            cx += cw

        for ri, row in enumerate(mom_rows):
            ry = ty + row_h * (ri + 1)
            _box(slide, tx, ry, total_w, row_h,
                 fill_color=WHITE if ri % 2 == 0 else ROW_ALT)

            change    = str(row.get("change", ""))
            direction = str(row.get("direction", ""))
            invert    = row.get("invert", False)

            if direction == "▲":
                chg_color = RED if invert else GREEN
            elif direction == "▼":
                chg_color = GREEN if invert else RED
            else:
                chg_color = DARK

            vals   = [row.get("metric",""), row.get("prev",""), row.get("curr",""), change, direction]
            colors = [DARK, DARK, DARK, chg_color, chg_color]
            aligns = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER,
                      PP_ALIGN.CENTER, PP_ALIGN.CENTER]
            bolds  = [True, False, False, True, False]

            cx = tx
            for val, col, al, bld, cw in zip(vals, colors, aligns, bolds, cw_list):
                _txt(slide, str(val), cx + Inches(0.04), ry + Inches(0.04),
                     cw - Inches(0.04), row_h - Inches(0.04),
                     size=9, bold=bld, color=col, align=al, font="Calibri")
                cx += cw

    _right_panel(slide,
                 "Key Takeaways",
                 _to_list(c.get("key_takeaways", [])),
                 "Next Steps",
                 _to_list(c.get("next_steps", [])))

    _sg_logo(slide)
    return slide


def _slide3(prs: Presentation, report: Dict):
    """Slide 3 — Design Center Scorecard."""
    slide = _blank_slide(prs)
    c = _get_slide_content(report, 3)
    month = report.get("report_month", "Monthly")

    stats_raw = c.get("summary_stats", {})
    _slide_header(slide, "Design Center Scorecard",
                  f"HubSpot Data  |  {month}  |  32 Locations")

    _stat_bar(slide, [
        ("Leads",   f"{stats_raw.get('total_leads', 0):,}"),
        ("Avg CPL", f"${stats_raw.get('avg_cpl', 0):,.2f}"),
        ("Visits",  f"{stats_raw.get('total_visits', 0):,}"),
        ("Quotes",  f"{stats_raw.get('total_quotes', 0):,}"),
    ], y=Inches(1.08))

    def _loc_table(section_title: str, rows: List[Dict], y_start, accent: RGBColor):
        # Accent bar + section title
        _box(slide, Inches(0.4), y_start, Inches(0.06), Inches(0.32), fill_color=accent)
        _txt(slide, section_title,
             Inches(0.54), y_start, Inches(7.5), Inches(0.32),
             size=10.5, bold=True, color=accent, font="Georgia")

        y       = y_start + Inches(0.35)
        tx      = Inches(0.4)
        cw_list = [Inches(2.0), Inches(0.9), Inches(0.9), Inches(1.05), Inches(0.9), Inches(1.05)]
        row_h   = Inches(0.275)
        total_w = sum(cw_list)
        headers = ["Location", "Leads", "Visits", "CPL", "Quotes", "Spend"]

        _box(slide, tx, y, total_w, row_h, fill_color=ORANGE)
        cx = tx
        for hdr, cw in zip(headers, cw_list):
            _txt(slide, hdr, cx + Inches(0.03), y + Inches(0.03),
                 cw, row_h,
                 size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
            cx += cw

        for ri, row in enumerate(rows[:6]):
            ry = y + row_h * (ri + 1)
            _box(slide, tx, ry, total_w, row_h,
                 fill_color=WHITE if ri % 2 == 0 else ROW_ALT)

            cpl_raw = row.get("cpl", 0)
            try:
                cpl_val   = float(str(cpl_raw).replace("$", "").replace(",", ""))
                cpl_color = GREEN if cpl_val < 100 else (RED if cpl_val > 175 else DARK)
                cpl_str   = f"${cpl_val:,.2f}"
            except (ValueError, TypeError):
                cpl_color = DARK
                cpl_str   = str(cpl_raw)

            visits_val = row.get("visits", 0)
            try:
                visit_num   = int(str(visits_val).replace(",", ""))
                visit_color = RED if visit_num <= 2 else DARK
            except (ValueError, TypeError):
                visit_color = DARK

            spend = row.get("spend", 0)
            spend_str = f"${spend:,.0f}" if isinstance(spend, (int, float)) else str(spend)

            vals   = [row.get("location",""), str(row.get("leads","")),
                      str(visits_val), cpl_str, str(row.get("quotes","")), spend_str]
            cols   = [DARK, DARK, visit_color, cpl_color, DARK, DARK]
            aligns = [PP_ALIGN.LEFT] + [PP_ALIGN.CENTER] * 5
            bolds  = [True] + [False] * 5

            cx = tx
            for val, col, al, bld, cw in zip(vals, cols, aligns, bolds, cw_list):
                _txt(slide, str(val), cx + Inches(0.03), ry + Inches(0.03),
                     cw, row_h - Inches(0.03),
                     size=8.5, bold=bld, color=col, align=al, font="Calibri")
                cx += cw

        return y + row_h * (min(len(rows), 6) + 1) + Inches(0.04)

    y_pos = Inches(2.04)
    y_pos = _loc_table("Top Performers — Leads & Visits",
                       c.get("top_performers", []), y_pos, GREEN)
    y_pos += Inches(0.08)
    _loc_table("Needs Attention — Low Volume or High CPL",
               c.get("needs_attention", []), y_pos, RED)

    _right_panel(slide,
                 "Key Insights",
                 _to_list(c.get("key_insights", [])),
                 "Focus Areas",
                 _to_list(c.get("focus_areas", [])))

    _sg_logo(slide)
    return slide


def _slide4(prs: Presentation, report: Dict):
    """Slide 4 — Attribution & Data Integrity."""
    slide = _blank_slide(prs)
    c = _get_slide_content(report, 4)

    _txt(slide, "Attribution & Data Integrity",
         Inches(0.4), Inches(0.15), RP - Inches(0.5), Inches(0.65),
         size=28, bold=True, color=DARK, font="Georgia")

    sync      = c.get("hubspot_sync", {})
    pmax      = c.get("pmax_status", "Monitoring")
    pixel     = c.get("meta_pixel_status", "Healthy")
    lead_sc   = c.get("lead_scoring_status", "In Progress")
    acc_table = c.get("accuracy_table", [])

    def _card(y_start, title, status, platforms=None, body=None, table=None):
        extra = Inches(0.3) * len(platforms or [])
        if table:
            extra += Inches(0.88)
        if body:
            extra += Inches(0.52)
        card_h = Inches(0.68) + extra

        _box(slide, Inches(0.4), y_start, RP - Inches(0.5), card_h,
             fill_color=WHITE, line_color=GRAY_BDR, line_pt=0.5)
        _box(slide, Inches(0.4), y_start, Inches(0.06), card_h, fill_color=ORANGE)

        _txt(slide, title,
             Inches(0.56), y_start + Inches(0.1), Inches(4.6), Inches(0.36),
             size=12, bold=True, color=DARK, font="Georgia")
        sc = _status_color(status)
        _txt(slide, f"[{status}]",
             Inches(5.3), y_start + Inches(0.11), Inches(3.0), Inches(0.3),
             size=10.5, bold=True, color=sc, font="Calibri")

        y = y_start + Inches(0.5)

        if platforms:
            px = Inches(0.62)
            for plat, pstat in platforms.items():
                pc = _status_color(pstat)
                _txt(slide, f"{plat}:", px, y, Inches(1.1), Inches(0.26),
                     size=9, bold=True, color=DARK, font="Calibri")
                _txt(slide, pstat, px + Inches(1.05), y, Inches(1.5), Inches(0.26),
                     size=9, color=pc, font="Calibri")
                px += Inches(2.6)
            y += Inches(0.32)

        if table:
            tbl_x   = Inches(0.56)
            tbl_cws = [Inches(1.4), Inches(1.1), Inches(1.1), Inches(1.5), Inches(1.0), Inches(0.9)]
            tbl_hdrs = ["Metric", "In-Platform", "HubSpot", "Variance", "Accuracy", "Target: 85%"]
            tr_h    = Inches(0.26)
            tbl_w   = sum(tbl_cws)

            _box(slide, tbl_x, y, tbl_w, tr_h, fill_color=ORANGE)
            cx = tbl_x
            for hdr, cw in zip(tbl_hdrs, tbl_cws):
                _txt(slide, hdr, cx + Inches(0.02), y + Inches(0.03), cw, tr_h,
                     size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
                cx += cw

            for ri, trow in enumerate(table):
                ry  = y + tr_h * (ri + 1)
                _box(slide, tbl_x, ry, tbl_w, tr_h,
                     fill_color=WHITE if ri % 2 == 0 else ROW_ALT)
                on_tgt  = trow.get("on_target")
                tgt_txt = "✓" if on_tgt is True else ("✗" if on_tgt is False else "—")
                tgt_col = GREEN if on_tgt is True else (RED if on_tgt is False else DARK)
                acc_col = GREEN if on_tgt is True else (RED if on_tgt is False else DARK)
                row_vals = [trow.get("metric",""), trow.get("platform",""),
                            trow.get("hubspot",""), trow.get("variance",""),
                            trow.get("accuracy",""), tgt_txt]
                row_cols = [DARK, DARK, DARK, DARK, acc_col, tgt_col]
                cx = tbl_x
                for val, col, cw in zip(row_vals, row_cols, tbl_cws):
                    _txt(slide, str(val), cx + Inches(0.02), ry + Inches(0.03), cw, tr_h,
                         size=8.5, color=col, align=PP_ALIGN.CENTER, font="Calibri")
                    cx += cw
            y += tr_h * (len(table) + 1) + Inches(0.06)

        if body:
            _txt(slide, str(body),
                 Inches(0.56), y, RP - Inches(0.9), Inches(0.52),
                 size=9, color=GRAY_SUB, font="Calibri", wrap=True)

        return y_start + card_h + Inches(0.1)

    y = Inches(0.93)
    y = _card(y, "HubSpot Attribution Sync",
              sync.get("overall_status", "In Progress"),
              platforms={
                  "Google":    sync.get("google_status",    "In Progress"),
                  "Meta":      sync.get("meta_status",      "In Progress"),
                  "Microsoft": sync.get("microsoft_status", "In Progress"),
              },
              table=acc_table)
    y = _card(y, "PMax Attribution", pmax,
              body="With Google + HubSpot Sync QA'd, PMax is climbing out of the learning phase.")
    y = _card(y, "Meta Pixel Health", pixel,
              body="Meta Pixel is firing and optimised for high-intent conversions.")
    _card(y, "Lead Scoring", lead_sc,
          body="Once established in HubSpot, lead scoring will unlock enhanced conversions in Google.")

    _right_panel(slide, "Action Items",
                 _to_list(c.get("action_items", [])))

    _sg_logo(slide)
    return slide


def _slide5(prs: Presentation, report: Dict):
    """Slide 5 — Ad Development & Testing (Top Creatives)."""
    slide = _blank_slide(prs)
    c = _get_slide_content(report, 5)

    _txt(slide, "Ad Development & Testing",
         Inches(0.4), Inches(0.15), RP - Inches(0.5), Inches(0.65),
         size=28, bold=True, color=DARK, font="Georgia", wrap=False)
    _box(slide, Inches(0.4), Inches(0.75), Inches(7.0), Pt(1.5), fill_color=DARK)

    creatives = c.get("creatives", report.get("top_creatives", []))
    y = Inches(0.98)

    for creative in creatives[:4]:
        ad_name     = creative.get("ad_name", "Ad Creative")
        campaign    = creative.get("campaign_name", "")
        leads       = creative.get("leads", 0)
        spend       = creative.get("spend", 0)
        cpl         = creative.get("cpl", 0)
        link_clicks = creative.get("link_clicks", creative.get("clicks", 0))

        _txt(slide, campaign,
             Inches(0.4), y, Inches(8.7), Inches(0.24),
             size=9, italic=True, bold=True, color=DARK, font="Calibri")
        y += Inches(0.25)

        card_h = Inches(0.9)
        _box(slide, Inches(0.4), y, Inches(1.3), card_h, fill_color=ROW_ALT)
        _txt(slide, "Ad\nCreative",
             Inches(0.42), y + Inches(0.28), Inches(1.25), Inches(0.42),
             size=7.5, color=GRAY_SUB, align=PP_ALIGN.CENTER, font="Calibri")

        result      = creative.get("result_label", "Top Creative")
        result_color = GREEN if "winner" in result.lower() else ORANGE
        _txt(slide, result,
             Inches(1.82), y + Inches(0.05), Inches(3.2), Inches(0.35),
             size=12, bold=True, color=result_color, font="Georgia")

        metrics_str = (
            f"{leads:,} Leads  |  ${cpl:,.2f} CPL  |  "
            f"{link_clicks:,} Link Clicks  |  ${spend:,.0f} Spend"
        )
        _txt(slide, metrics_str,
             Inches(1.82), y + Inches(0.40), Inches(7.1), Inches(0.25),
             size=8.5, color=DARK, font="Calibri")

        note = creative.get("note", "")
        if note:
            _txt(slide, note,
                 Inches(1.82), y + Inches(0.62), Inches(7.1), Inches(0.28),
                 size=8.5, color=GRAY_SUB, font="Calibri")

        y += card_h + Inches(0.2)

    _right_panel(slide,
                 "Key Takeaways",
                 _to_list(c.get("key_takeaways", ["Exterior Shots tend to perform best."])),
                 "Next Steps",
                 _to_list(c.get("next_steps", ["Iterate on winning creative styles."])))

    _sg_logo(slide)
    return slide


def _slide6(prs: Presentation, report: Dict):
    """Slide 6 — Current Initiatives & Priority Updates."""
    slide = _blank_slide(prs)
    c = _get_slide_content(report, 6)

    _oval(slide, Inches(9.4), Inches(0.3), Inches(5.2), Inches(5.2), PEACH)

    _txt(slide, "Current Initiatives & Priority Updates",
         Inches(0.4), Inches(0.18), Inches(9.0), Inches(0.65),
         size=26, bold=True, color=ORANGE, font="Georgia")

    raw = c.get("initiatives", [])
    parsed: List[Dict] = []
    for item in raw:
        item = str(item).strip()
        if not item:
            continue
        if ":" in item:
            parts   = item.split(":", 1)
            title   = parts[0].strip()
            bullets = [b.strip() for b in parts[1].split("|") if b.strip()]
        else:
            title   = item
            bullets = []
        parsed.append({"title": title, "bullets": bullets})

    y = Inches(1.02)
    for i, initiative in enumerate(parsed[:5]):
        title   = initiative.get("title", "")
        bullets = initiative.get("bullets", [])
        card_h  = Inches(0.62) + Inches(0.26) * len(bullets)

        _box(slide, Inches(0.4), y, Inches(8.7), card_h,
             fill_color=WHITE, line_color=GRAY_BDR, line_pt=0.5)
        _box(slide, Inches(0.4), y, Inches(0.05), card_h, fill_color=ORANGE)
        _oval(slide, Inches(0.56), y + Inches(0.11), Inches(0.42), Inches(0.42), ORANGE)
        _txt(slide, str(i + 1),
             Inches(0.56), y + Inches(0.09), Inches(0.42), Inches(0.42),
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
        _txt(slide, title,
             Inches(1.1), y + Inches(0.12), Inches(7.4), Inches(0.36),
             size=12, bold=True, color=DARK, font="Georgia")

        by = y + Inches(0.48)
        for bullet in bullets:
            _txt(slide, f"●  {bullet}",
                 Inches(1.3), by, Inches(7.1), Inches(0.26),
                 size=9, color=DARK, font="Calibri")
            by += Inches(0.26)

        y += card_h + Inches(0.1)

    _sg_logo(slide)
    return slide


def _slide7(prs: Presentation, report: Dict):
    """Slide 7 — Strategic Recommendations."""
    slide = _blank_slide(prs)
    c = _get_slide_content(report, 7)

    _txt(slide, "Strategic Recommendations",
         Inches(0.4), Inches(0.15), RP - Inches(0.5), Inches(0.65),
         size=28, bold=True, color=DARK, font="Georgia")

    recs = c.get("recommendations", [])
    y    = Inches(0.94)

    for rec in recs[:4]:
        title   = rec.get("title", "")
        body    = rec.get("body", "")
        bullets = _to_list(body) if body else []

        card_h  = Inches(0.52) + Inches(0.3) * max(len(bullets), 1)

        _box(slide, Inches(0.4), y, RP - Inches(0.5), card_h,
             fill_color=WHITE, line_color=GRAY_BDR, line_pt=0.5)
        _box(slide, Inches(0.4), y, Inches(0.06), card_h, fill_color=ORANGE)
        _txt(slide, title,
             Inches(0.58), y + Inches(0.09), RP - Inches(0.92), Inches(0.36),
             size=12, bold=True, color=DARK, font="Georgia")

        by = y + Inches(0.45)
        for bullet in bullets:
            _txt(slide, f"●  {bullet}",
                 Inches(0.66), by, RP - Inches(1.0), Inches(0.3),
                 size=9, color=DARK, font="Calibri", wrap=True)
            by += Inches(0.3)

        y += card_h + Inches(0.12)

    _right_panel(slide, "What's Next",
                 _to_list(c.get("whats_next", "")))

    _sg_logo(slide)
    return slide


# ─── Public API ───────────────────────────────────────────────────────────────

_BUILDERS = [_slide1, _slide2, _slide3, _slide4, _slide5, _slide6, _slide7]


def build_pptx(report: Dict[str, Any]) -> bytes:
    """
    Build a 7-slide branded PPTX from a MonthlySlidesResponse dict.
    Returns raw bytes suitable for streaming or uploading to Google Drive.
    """
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    for i, builder in enumerate(_BUILDERS):
        try:
            builder(prs, report)
        except Exception as exc:
            import traceback
            print(f"[pptx_builder] Error on slide {i + 1}: {exc}")
            traceback.print_exc()
            err_slide = _blank_slide(prs)
            _txt(err_slide, f"Slide {i + 1} — Generation Error",
                 Inches(1), Inches(3), Inches(11), Inches(1),
                 size=16, bold=True, color=RED)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
